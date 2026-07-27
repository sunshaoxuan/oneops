import argparse
import io
import json
import os
from pathlib import Path
import subprocess
import sys
import tempfile


def normalized_text(value):
    lines = [line.rstrip() for line in str(value or "").replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    result = "\n".join(lines).strip()
    while "\n\n\n" in result:
        result = result.replace("\n\n\n", "\n\n")
    return result


def decode_text(data):
    for encoding in ("utf-8-sig", "cp932", "shift_jis", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def ocr_pdf(data, page_count):
    import pypdfium2 as pdfium

    maximum_ocr_pages = 60
    rendered_pages = min(page_count, maximum_ocr_pages)
    script = Path(__file__).with_name("inquiry_attachment_ocr.ps1")
    powershell = os.environ.get(
        "OPS_INQUIRY_ATTACHMENT_POWERSHELL",
        r"C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe",
    )
    with tempfile.TemporaryDirectory(prefix="oneops-inquiry-ocr-") as directory:
        pdf = pdfium.PdfDocument(data)
        try:
            for index in range(rendered_pages):
                page = pdf[index]
                try:
                    bitmap = page.render(scale=2)
                    image = bitmap.to_pil()
                    image.save(
                        Path(directory) / f"page-{index + 1:04d}.png",
                        "PNG",
                    )
                finally:
                    page.close()
        finally:
            pdf.close()
        completed = subprocess.run(
            [
                powershell,
                "-NoLogo",
                "-NoProfile",
                "-NonInteractive",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(script),
                "-Directory",
                directory,
            ],
            capture_output=True,
            check=False,
            encoding="utf-8-sig",
            errors="replace",
            timeout=100,
        )
    output = completed.stdout.strip()
    first_object = output.find("{")
    last_object = output.rfind("}")
    if completed.returncode != 0 or first_object < 0 or last_object < first_object:
        raise RuntimeError(
            completed.stderr.strip()[:500] or
            "Windows OCR returned an invalid response."
        )
    result = json.loads(output[first_object:last_object + 1])
    if result.get("status") != "PARSED":
        raise RuntimeError(
            result.get("message") or "Windows OCR failed."
        )
    parts = []
    for index, page in enumerate(result.get("pages", []), start=1):
        text = normalized_text(page.get("text", ""))
        if text:
            parts.append(f"[Page {index}]\n{text}")
    return {
        "text": "\n\n".join(parts),
        "parser": "windows-ocr-ja-JP",
        "pageCount": page_count,
        "truncated": page_count > rendered_pages,
    }


def parse_pdf(data):
    import pdfplumber
    from pypdf import PdfReader

    pages = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            value = normalized_text(page.extract_text(layout=True) or "")
            if value:
                pages.append(f"[Page {index}]\n{value}")
        page_count = len(pdf.pages)
    if not pages:
        reader = PdfReader(io.BytesIO(data))
        for index, page in enumerate(reader.pages, start=1):
            value = normalized_text(page.extract_text() or "")
            if value:
                pages.append(f"[Page {index}]\n{value}")
        page_count = len(reader.pages)
    if not pages and page_count:
        return ocr_pdf(data, page_count)
    return {
        "text": "\n\n".join(pages),
        "parser": "pdfplumber+pypdf",
        "pageCount": page_count,
        "truncated": False,
    }


def parse_docx(data):
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return {
        "text": "\n".join(parts),
        "parser": "python-docx",
    }


def parse_xlsx(data):
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for worksheet in workbook.worksheets:
        parts.append(f"[Sheet: {worksheet.title}]")
        for row in worksheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                parts.append("\t".join(values))
    return {
        "text": "\n".join(parts),
        "parser": "openpyxl",
        "sheetCount": len(workbook.worksheets),
    }


def parse_pptx(data):
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        values = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                values.append(shape.text)
        if values:
            parts.append(f"[Slide {index}]\n" + "\n".join(values))
    return {
        "text": "\n\n".join(parts),
        "parser": "python-pptx",
        "slideCount": len(presentation.slides),
    }


def parse_attachment(data, name, content_type):
    extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    if extension == "pdf" or content_type == "application/pdf":
        return parse_pdf(data)
    if extension == "docx" or "wordprocessingml" in content_type:
        return parse_docx(data)
    if extension == "xlsx" or "spreadsheetml" in content_type:
        return parse_xlsx(data)
    if extension == "pptx" or "presentationml" in content_type:
        return parse_pptx(data)
    if extension in {
        "txt",
        "csv",
        "tsv",
        "log",
        "md",
        "json",
        "xml",
        "html",
        "htm",
        "yaml",
        "yml",
        "sql",
    } or content_type.startswith("text/"):
        return {
            "text": decode_text(data),
            "parser": "text-decoder",
        }
    return {
        "status": "UNSUPPORTED",
        "text": "",
        "parser": "",
        "errorCode": "UNSUPPORTED_FORMAT",
        "message": f"Unsupported attachment format: {extension or content_type or 'unknown'}",
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--name", required=True)
    parser.add_argument("--content-type", default="")
    args = parser.parse_args()
    try:
        result = parse_attachment(
            sys.stdin.buffer.read(),
            args.name,
            args.content_type.lower().split(";", 1)[0].strip(),
        )
        result["text"] = normalized_text(result.get("text", ""))
        if "status" not in result:
            result["status"] = "PARSED" if result["text"] else "EMPTY"
        if result["status"] == "EMPTY":
            result["errorCode"] = "NO_EXTRACTABLE_TEXT"
            result["message"] = "The attachment contains no recognizable text."
    except Exception as error:
        result = {
            "status": "FAILED",
            "text": "",
            "parser": "",
            "errorCode": error.__class__.__name__,
            "message": str(error)[:500],
        }
    sys.stdout.write(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
